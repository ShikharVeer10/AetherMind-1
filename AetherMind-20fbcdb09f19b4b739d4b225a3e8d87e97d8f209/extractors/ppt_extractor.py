from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from models.document_model import (
    DocumentModel,
    DocumentElementModel,
    SlideModel,
    PositionModel,
    StyleModel,
    ParagraphModel,
    RunModel,
)
from services.table_service import TableService
from models.document_model import (
    ChartUnderstandingModel,
    ChartSeriesModel,
    ChartAxisModel,
)


class PPTExtractor:
    def __init__(self, pptx_file_path: str):
        self.pptx_file_path = Path(pptx_file_path)
        self.presentation = Presentation(pptx_file_path)
        self.table_service = TableService()

    def extract_document(self) -> DocumentModel:
        extracted_slides = []
        for slide_index, slide in enumerate(self.presentation.slides):
            extracted_slide = self.extract_slide(
                slide=slide, slide_number=slide_index + 1
            )
            extracted_slides.append(extracted_slide)

        # Extract presentation-level metadata
        pres_meta = self._extract_presentation_metadata()

        return DocumentModel(
            document_name=Path(self.pptx_file_path).name,
            document_type="ppt",
            total_slides=len(extracted_slides),
            slides=extracted_slides,
            presentation_metadata=pres_meta,
        )
    

    def _extract_presentation_metadata(self) -> dict:
        """Extract top-level presentation metadata (author, dimensions, etc.)"""
        meta: dict = {}
        try:
            core = self.presentation.core_properties
            if core.author:
                meta["author"] = core.author
            if core.title:
                meta["presentation_title"] = core.title
            if core.subject:
                meta["subject"] = core.subject
            if core.created:
                meta["created"] = str(core.created)
            if core.modified:
                meta["modified"] = str(core.modified)
        except Exception:
            pass
        try:
            meta["slide_width"] = float(self.presentation.slide_width)
            meta["slide_height"] = float(self.presentation.slide_height)
        except Exception:
            pass
        return meta

    def _extract_list_properties_from_xml(self, paragraph) -> dict:
        list_props = {}
        try:
            pPr = paragraph._p.pPr
            if pPr is not None:
                buChar = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                if buChar is not None:
                    list_props["bullet_character"] = buChar.get("char")
                    list_props["list_type"] = "bulleted"
                else:
                    buAutoNum = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    if buAutoNum is not None:
                        list_props["list_type"] = "numbered"
                        list_props["number_type"] = buAutoNum.get("type", "arabicPeriod")
        except Exception:
            pass
        return list_props

    def _extract_slide_background(self, slide) -> dict:
        bg_info = {"type": "solid", "color": "#ffffff", "opacity": 1.0}
        try:
            background = slide.background
            if background and background.fill:
                fill = background.fill
                fill_type = getattr(fill, "type", None)
                if fill_type == 1: # Solid
                    bg_info["type"] = "solid"
                    if fill.fore_color:
                        bg_info["color"] = self._get_safe_color_hex(fill.fore_color) or "#ffffff"
                    bg_info["opacity"] = 1.0 - float(getattr(fill, "transparency", 0.0))
                elif fill_type == 3: # Gradient
                    bg_info["type"] = "gradient"
                    bg_info["gradient_details"] = {}
                    try:
                        stops = []
                        for stop in fill.gradient_stops:
                            stops.append({
                                "color": self._get_safe_color_hex(stop.color),
                                "position": stop.position
                            })
                        bg_info["gradient_details"]["stops"] = stops
                    except Exception:
                        pass
                    if fill.fore_color:
                        bg_info["color"] = self._get_safe_color_hex(fill.fore_color) or "#ffffff"
                elif fill_type == 6: # Picture
                    bg_info["type"] = "picture"
                    bg_info["role"] = "background"
                    bg_info["color"] = None
        except Exception:
            pass
        return bg_info

    def extract_slide(self, slide, slide_number: int) -> SlideModel:
        extracted_elements = []
        slide_title: Optional[str] = None

        # 1. Try to get the real title from placeholders
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = slide.shapes.title.text_frame.text.strip() or None
        except Exception:
            pass

        for index, shape in enumerate(slide.shapes):
            elements = self._extract_shape_recursive(
                shape=shape,
                slide_number=slide_number,
                shape_index=index,
                z_order=index,
                prefix=f"slide_{slide_number}",
            )
            for element in elements:
                extracted_elements.append(element)
        
        z_order = len(slide.shapes)
        # 2. Fallback: if no placeholder title, use the first element text
        if slide_title is None:
            for element in extracted_elements:
                if element.text:
                    slide_title = element.text.strip()
                    break

        # Extract slide background color
        bg_info = self._extract_slide_background(slide)
        slide_bg_color = bg_info.get("color") or "#ffffff"

        # 3. Visual Table Detection Fallback
        from services.flexible_table_detector import FlexibleTableDetector
        detector = FlexibleTableDetector()
        visual_tables = detector.detect_visual_tables(extracted_elements)
        
        detected_tables_metadata = []
        consumed_element_ids = set()
        
        for vt_idx, visual_table in enumerate(visual_tables):
            raw_table_content = visual_table.get("rows", [])
            raw_table_styles = visual_table.get("styles", [])
            if not raw_table_content:
                continue

            consumed_element_ids.update(visual_table.get("consumed_ids", []))

            table_md = self.table_service.to_markdown(raw_table_content)
            table_structure = self.table_service.analyze_structure(raw_table_content)
            table_structure["merged_cells"] = visual_table.get("merged_cells", [])
            
            table_semantic_interpretation = self.table_service.generate_semantic_context(raw_table_content)
            table_render_model = self.table_service.build_render_model(raw_table_content, table_structure)
            bbox = visual_table.get("bbox", {"x": 0, "y": 0, "width": 0, "height": 0})

            table_reconstruction = self.table_service.build_reconstruction_payload(
                table_id=f"slide_{slide_number}_vtable_{vt_idx}",
                raw_table_content=raw_table_content,
                table_structure=table_structure,
                table_render_model=table_render_model,
                table_semantics=table_semantic_interpretation,
                is_visual=True,
                table_geometry=bbox,
                raw_table_styles=raw_table_styles,
                row_heights=visual_table.get("row_heights"),
                column_widths=visual_table.get("column_widths")
            )

            table_element = DocumentElementModel(
                element_id=f"slide_{slide_number}_vtable_{vt_idx}",
                element_type="table",
                text=table_md,
                paragraphs=[],
                position=PositionModel(
                    x=bbox["x"],
                    y=bbox["y"],
                    width=bbox["width"],
                    height=bbox["height"]
                ),
                table_markdown=table_md,
                raw_table_content=raw_table_content,
                table_structure=table_structure,
                table_render_model=table_render_model,
                table_semantic_interpretation=table_semantic_interpretation,
                table_reconstruction=table_reconstruction,
                table_merged_cells=visual_table.get("merged_cells", []),
                metadata={"name": f"Visual Table {vt_idx}", "z_order": z_order}
            )
            extracted_elements.append(table_element)
            detected_tables_metadata.append({
                "table_type": "visual_table",
                "rows": len(raw_table_content),
                "content": raw_table_content
            })
            z_order += 1

        # 4. Filter out consumed elements from the main list
        final_elements = [e for e in extracted_elements if e.element_id not in consumed_element_ids]

        slide_model = SlideModel(
            slide_number=slide_number,
            title=slide_title,
            elements=final_elements,
            background_color=slide_bg_color,
            detected_tables=detected_tables_metadata
        )
        if not slide_model.metadata:
            slide_model.metadata = {}
        slide_model.metadata["background"] = bg_info
        return slide_model

    def _extract_shape_recursive(
        self,
        shape,
        slide_number: int,
        shape_index: int,
        prefix: str,
        z_order: int = 0,
    ) -> List[DocumentElementModel]:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_id = f"{prefix}_group_{shape_index}"
            x, y, w, h = self._safe_position(shape)
            group_element = DocumentElementModel(
                element_id=group_id,
                element_type="group",
                text="",
                paragraphs=[],
                position=PositionModel(x=x, y=y, width=w, height=h),
                shape_type="group",
                metadata={
                    "name": shape.name,
                    "z_order": z_order,
                    "rotation": float(getattr(shape, "rotation", 0) or 0),
                    "children": []
                }
            )
            child_elements = []
            group_shapes = getattr(shape, "shapes", [])
            for child_idx, child_shape in enumerate(group_shapes):
                children_extracted = self._extract_shape_recursive(
                    shape=child_shape,
                    slide_number=slide_number,
                    shape_index=child_idx,
                    z_order=z_order + child_idx + 1,
                    prefix=group_id,
                )
                for c_elem in children_extracted:
                    if not c_elem.metadata.get("parent"):
                        c_elem.metadata["parent"] = group_id
                        c_elem.metadata["group_id"] = group_id
                        c_elem.parent = group_id
                        c_elem.group_id = group_id
                    group_element.metadata["children"].append(c_elem.element_id)
                    group_element.children.append(c_elem.element_id)
                    child_elements.append(c_elem)
            return [group_element] + child_elements
            
        element = self._extract_single_shape(
            shape=shape,
            element_id=f"{prefix}_shape_{shape_index}",
            z_order=z_order,
        )
        return [element] if element is not None else []

    def _extract_single_shape(
        self, shape, element_id: str, z_order: int = 0
    ) -> Optional[DocumentElementModel]:
        paragraphs: List[ParagraphModel] = []
        full_text: Optional[str] = None

        if shape.has_text_frame:
            paragraph_texts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if not para_text:
                    continue

                runs = []
                for run in paragraph.runs:
                    run_text = run.text
                    if not run_text:
                        continue
                    run_model = RunModel(
                        text=run_text,
                        bold=bool(run.font.bold),
                        italic=bool(run.font.italic),
                        font_size=(
                            float(run.font.size.pt) if run.font.size else None
                        ),
                        font_name=run.font.name,
                        font_color=self._get_safe_color_hex(run.font.color),
                        underline=bool(run.font.underline),
                        font_weight="bold" if run.font.bold else "normal",
                        letter_spacing=None,
                    )
                    runs.append(run_model)

                # Extract paragraph alignment
                alignment_name = None
                if paragraph.alignment is not None:
                    try:
                        alignment_name = paragraph.alignment.name
                    except Exception:
                        alignment_name = str(paragraph.alignment)

                # Spacing & Indentation
                line_spacing = None
                if paragraph.line_spacing is not None:
                    try:
                        if hasattr(paragraph.line_spacing, "pt"):
                            line_spacing = float(paragraph.line_spacing.pt)
                        else:
                            line_spacing = float(paragraph.line_spacing)
                    except Exception:
                        pass
                
                space_before = None
                if paragraph.space_before is not None:
                    try:
                        space_before = float(paragraph.space_before.pt)
                    except Exception:
                        pass

                space_after = None
                if paragraph.space_after is not None:
                    try:
                        space_after = float(paragraph.space_after.pt)
                    except Exception:
                        pass

                indentation = None
                try:
                    if hasattr(paragraph, "left_indent") and paragraph.left_indent is not None:
                        indentation = float(paragraph.left_indent.pt)
                except Exception:
                    pass

                # List / Bullet Info
                list_props = self._extract_list_properties_from_xml(paragraph)
                list_type = list_props.get("list_type")
                bullet_char = list_props.get("bullet_character")
                number_val = list_props.get("number_type") if list_type == "numbered" else None

                paragraphs.append(
                    ParagraphModel(
                        level=paragraph.level,
                        text=para_text,
                        runs=runs,
                        alignment=alignment_name,
                        line_spacing=line_spacing,
                        paragraph_spacing_before=space_before,
                        paragraph_spacing_after=space_after,
                        indentation=indentation,
                        list_type=list_type,
                        bullet_character=bullet_char,
                        number=number_val,
                    )
                )

                paragraph_texts.append(para_text)

            if paragraph_texts:
                full_text = "\n".join(paragraph_texts)

        x, y, w, h = self._safe_position(shape)
        position = PositionModel(
            x=x,
            y=y,
            width=w,
            height=h,
        )

        style = self._extract_text_style(shape)
        element_type = self._get_shape_type(shape)
        print(
            "SHAPE:",
            shape.name,
            "TYPE:",
            shape.shape_type,
            "ELEMENT:",
            element_type
        )       
        print("=" * 80)
        print("SHAPE NAME:", shape.name)
        print("SHAPE TYPE:", shape.shape_type)
        print("ELEMENT TYPE:", element_type)
        print("=" * 80)
        metadata = self._extract_shape_metadata(shape, element_type)
        metadata["z_order"] = z_order

        # Extract rich shape formatting properties
        shape_geom_type = "rect"
        try:
            if hasattr(shape, "auto_shape_type") and shape.auto_shape_type is not None:
                auto_name = shape.auto_shape_type.name.lower()
                if "oval" in auto_name or "circle" in auto_name:
                    shape_geom_type = "circle"
                elif "round" in auto_name:
                    shape_geom_type = "rounded_rectangle"
                else:
                    shape_geom_type = "rectangle"
        except (ValueError, AttributeError):
            pass
        if element_type in {"connector", "line"}:
            shape_geom_type = "line"

        fill_color = None
        fill_type = "solid"
        opacity = 1.0
        gradient = None
        try:
            if hasattr(shape, "fill") and shape.fill:
                ft = getattr(shape.fill, "type", None)
                if ft == 1: # Solid
                    if shape.fill.fore_color:
                        fill_color = self._get_safe_color_hex(shape.fill.fore_color)
                elif ft == 3: # Gradient
                    fill_type = "gradient"
                    gradient = {}
                    try:
                        stops = []
                        for stop in shape.fill.gradient_stops:
                            stops.append({
                                "color": self._get_safe_color_hex(stop.color),
                                "position": stop.position
                            })
                        gradient["stops"] = stops
                    except Exception:
                        pass
                    if shape.fill.fore_color:
                        fill_color = self._get_safe_color_hex(shape.fill.fore_color)
                elif ft == 6: # Picture
                    fill_type = "image"
                
                try:
                    opacity = 1.0 - float(getattr(shape.fill, "transparency", 0.0))
                except Exception:
                    pass
        except Exception:
            pass

        border_color = None
        border_thickness = None
        try:
            if hasattr(shape, "line") and shape.line:
                border_color = self._get_safe_color_hex(shape.line.color)
                if shape.line.width:
                    border_thickness = float(shape.line.width.pt)
        except Exception:
            pass

        border_radius = None
        try:
            if hasattr(shape, "adjustments") and shape.adjustments and len(shape.adjustments) > 0:
                border_radius = float(shape.adjustments[0])
        except Exception:
            pass

        shadow = None
        try:
            if hasattr(shape, "shadow") and shape.shadow:
                if shape.shadow.inherit is False:
                    shadow = {
                        "color": self._get_safe_color_hex(shape.shadow.color) if hasattr(shape.shadow, "color") else None,
                        "blur_radius": float(shape.shadow.blur_radius.pt) if hasattr(shape.shadow, "blur_radius") and shape.shadow.blur_radius else None,
                        "distance": float(shape.shadow.distance.pt) if hasattr(shape.shadow, "distance") and shape.shadow.distance else None,
                        "angle": float(shape.shadow.direction) if hasattr(shape.shadow, "direction") else None,
                    }
        except Exception:
            pass

        crop = None
        mask = None
        caption = None
        role = None
        if element_type == "image":
            crop = {
                "left": float(getattr(shape, "crop_left", 0.0)),
                "top": float(getattr(shape, "crop_top", 0.0)),
                "right": float(getattr(shape, "crop_right", 0.0)),
                "bottom": float(getattr(shape, "crop_bottom", 0.0)),
            }
            mask = "none"
            caption = getattr(shape, "description", "") or ""
            role = "illustration"
            w_pt = float(shape.width) / 12700.0 if hasattr(shape, "width") else 0.0
            h_pt = float(shape.height) / 12700.0 if hasattr(shape, "height") else 0.0
            if w_pt < 48 and h_pt < 48:
                role = "icon"
            try:
                pres_width = float(self.presentation.slide_width)
                pres_height = float(self.presentation.slide_height)
                if pres_width > 0 and pres_height > 0:
                    area_ratio = (float(shape.width) * float(shape.height)) / (pres_width * pres_height)
                    if area_ratio > 0.8:
                        role = "background"
            except Exception:
                pass

        metadata.update({
            "fill_color": fill_color,
            "border_color": border_color,
            "border_thickness": border_thickness,
            "border_radius": border_radius,
            "shadow": shadow,
            "gradient": gradient,
            "opacity": opacity,
            "shape_type": shape_geom_type,
        })
        if element_type == "image":
            metadata.update({
                "crop": crop,
                "mask": mask,
                "caption": caption,
                "role": role,
            })

        if style:
            style.background_color = fill_color
            style.border_color = border_color
            style.border_thickness = border_thickness
            style.border_radius = border_radius
            style.opacity = opacity
            style.shadow = shadow
            style.gradient = gradient
            style.underline = any(r.underline for p in paragraphs for r in p.runs)

        table_md = None
        raw_table_content = None
        table_structure = None
        table_render_model = None
        table_semantic_interpretation = None
        table_visual_metadata = None
        table_reconstruction = None

        if element_type == "table":
            raw_table_content = self.extract_table_as_list(shape)
            raw_table_styles = self._extract_table_styles(shape)
            table_md = self.extract_table_as_markdown(shape)
            # Ensure text attribute contains the table markdown for verbatim extraction
            full_text = table_md

            table_structure = self.table_service.analyze_structure(raw_table_content)
            table_semantic_interpretation = self.table_service.generate_semantic_context(raw_table_content)
            table_render_model = self.table_service.build_render_model(raw_table_content, table_structure)
            table_visual_metadata = self.extract_table_visual_metadata(shape)

            table_reconstruction = self.table_service.build_reconstruction_payload(
                table_id=element_id,
                raw_table_content=raw_table_content,
                table_structure=table_structure,
                table_render_model=table_render_model,
                table_semantics=table_semantic_interpretation,
                is_visual=False,
                table_geometry={"x": x, "y": y, "width": w, "height": h},
                raw_table_styles=raw_table_styles,
                row_heights=table_visual_metadata.get("row_heights"),
                column_widths=table_visual_metadata.get("column_widths")
            )
        else:
            table_md = self.extract_table_as_markdown(shape)

        elem = DocumentElementModel(
            element_id=element_id,
            element_type=element_type,
            text=full_text,
            paragraphs=paragraphs,
            position=position,
            style=style,
            shape_type=shape_geom_type,
            metadata=metadata,
            table_markdown=table_md,
            raw_table_content=raw_table_content,
            table_structure=table_structure,
            table_render_model=table_render_model,
            table_semantic_interpretation=table_semantic_interpretation,
            table_visual_metadata=table_visual_metadata,
            table_reconstruction=table_reconstruction,
            
            # Enrich direct fields
            border_radius=border_radius,
            shadow=shadow,
            gradient=gradient,
            opacity=opacity,
            underline=any(r.underline for p in paragraphs for r in p.runs),
            caption=caption,
            role=role,
            crop=crop,
            mask=mask,
        )

        first_list_para = next((p for p in paragraphs if p.list_type), None)
        if first_list_para:
            elem.bullet_character = first_list_para.bullet_character
            elem.indentation = first_list_para.indentation
            elem.level = first_list_para.level
            elem.number = first_list_para.number

        return elem

    def _extract_cell_padding(self, cell) -> dict:
        padding = {}
        try:
            if cell.margin_left is not None:
                padding["padding_left"] = float(cell.margin_left.pt)
            if cell.margin_top is not None:
                padding["padding_top"] = float(cell.margin_top.pt)
            if cell.margin_right is not None:
                padding["padding_right"] = float(cell.margin_right.pt)
            if cell.margin_bottom is not None:
                padding["padding_bottom"] = float(cell.margin_bottom.pt)
        except Exception:
            pass
        return padding

    def _extract_cell_borders(self, cell) -> dict:
        borders = {}
        try:
            tcPr = cell._tc.get_or_add_tcPr()
            for side, xml_tag in [("top", "lnT"), ("bottom", "lnB"), ("left", "lnL"), ("right", "lnR")]:
                ln = tcPr.find(f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{xml_tag}')
                if ln is not None:
                    width = float(ln.get("w", 12700)) / 12700.0 # w is in EMUs (1 pt = 12700 EMUs)
                    color = "#000000"
                    solidFill = ln.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    if solidFill is not None:
                        srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        if srgbClr is not None:
                            val = srgbClr.get("val")
                            if val:
                                color = f"#{val.lower()}"
                    borders[f"border_{side}"] = {"width": width, "color": color, "style": "solid"}
                else:
                    borders[f"border_{side}"] = {"width": 0.0, "color": "#000000", "style": "none"}
        except Exception:
            pass
        return borders

    def _extract_table_styles(self, shape) -> List[List[StyleModel]]:
        """Extracts styling (bg color, font properties) for every cell in a native PPT table."""
        styles = []
        if not shape.has_table:
            return []
        
        for row in shape.table.rows:
            row_styles = []
            for cell in row.cells:
                # 1. Background color
                bg_color = None
                try:
                    if cell.fill and cell.fill.fore_color:
                        bg_color = self._get_safe_color_hex(cell.fill.fore_color)
                except Exception:
                    pass
                
                # 2. Font properties from first paragraph
                font_size = None
                font_name = None
                is_bold = False
                is_italic = False
                is_underline = False
                text_color = None
                try:
                    if cell.text_frame.paragraphs:
                        p = cell.text_frame.paragraphs[0]
                        if p.runs:
                            f = p.runs[0].font
                            font_size = float(f.size.pt) if f.size else None
                            font_name = f.name
                            is_bold = bool(f.bold)
                            is_italic = bool(f.italic)
                            is_underline = bool(f.underline)
                            text_color = self._get_safe_color_hex(f.color)
                except Exception:
                    pass
                
                # 3. Vertical alignment
                v_align = "center"
                try:
                    if cell.vertical_anchor is not None:
                        va_name = str(cell.vertical_anchor).lower()
                        if "top" in va_name:
                            v_align = "top"
                        elif "bottom" in va_name:
                            v_align = "bottom"
                except Exception:
                    pass

                padding = self._extract_cell_padding(cell)
                borders = self._extract_cell_borders(cell)

                row_styles.append(StyleModel(
                    background_color=bg_color,
                    font_size=font_size,
                    font_name=font_name,
                    bold=is_bold,
                    italic=is_italic,
                    underline=is_underline,
                    text_color=text_color,
                    vertical_alignment=v_align,
                    padding=padding,
                    border_top=borders.get("border_top"),
                    border_bottom=borders.get("border_bottom"),
                    border_left=borders.get("border_left"),
                    border_right=borders.get("border_right")
                ))
            styles.append(row_styles)
        return styles
        

    def _get_shape_type(self, shape) -> str:
        st = shape.shape_type

        if st == MSO_SHAPE_TYPE.TEXT_BOX:
            return "text_box"

        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            name = shape.name.lower()
            # Detect arrows from shape name
            if "arrow" in name:
                return "arrow"
            # Detect process/decision shapes common in flowcharts
            auto_shape_type = getattr(shape, 'auto_shape_type', None)
            if auto_shape_type is not None:
                auto_name = str(auto_shape_type).lower()
                if 'arrow' in auto_name:
                    return "arrow"
                if 'diamond' in auto_name or 'decision' in auto_name:
                    return "shape"  # decision diamond
                if 'chevron' in auto_name or 'pentagon' in auto_name:
                    return "shape"  # process chevron
            return "shape"

        if st == MSO_SHAPE_TYPE.GROUP:
            return "group"

        if st == MSO_SHAPE_TYPE.TABLE:
            return "table"

        if st == MSO_SHAPE_TYPE.PICTURE:
            return "image"

        if st == MSO_SHAPE_TYPE.CHART:
            return "chart"

        if st == MSO_SHAPE_TYPE.FREEFORM:
            return "freeform"

        if st == MSO_SHAPE_TYPE.PLACEHOLDER:
            return "placeholder"

        connector_type = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
        if connector_type is not None and st == connector_type:
            return "connector"
        if st == MSO_SHAPE_TYPE.LINE:
            return "connector"

        return "unknown"


    def _extract_shape_metadata(self, shape, element_type: str) -> dict:
        metadata: dict = {}
        metadata["name"] = getattr(shape, "name", "")
        try:
            metadata["rotation"] = float(shape.rotation or 0)
        except Exception:
            metadata["rotation"] = 0
        try:
            metadata["visible"] = bool(getattr(shape, "visible", True))
        except Exception:
            metadata["visible"] = True
        try:
            metadata["is_placeholder"] = bool(getattr(shape, "is_placeholder", False))
        except Exception:
            metadata["is_placeholder"] = False

        if element_type == "table":
            metadata["table_data"] = self.extract_table_data(shape)
        if element_type == "chart":
            metadata["chart_data"] = self._extract_pptx_chart_data(shape)
        if element_type == "image":
            try:
                metadata["__image_bytes"] = shape.image.blob
            except Exception:
                try:
                    # Fallback to direct relationship blob access for custom content type images (e.g. image/jpg)
                    blip = shape._element.blipFill.blip
                    embed_key = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                    link_key = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link'
                    rId = blip.get(embed_key) or blip.get(link_key)
                    rel = shape.part.rels[rId]
                    target_part = rel.target_part
                    if hasattr(target_part, "blob"):
                        metadata["__image_bytes"] = target_part.blob
                    elif hasattr(target_part, "_blob"):
                        metadata["__image_bytes"] = target_part._blob
                    else:
                        metadata["__image_bytes"] = None
                except Exception:
                    metadata["__image_bytes"] = None

        if shape.has_text_frame:
            metadata["bullet_hierarchy"] = self._extract_bullet_hierarchy(shape)
        if element_type == "connector":
            metadata["connector_endpoints"] = (
                self._extract_connector_endpoints(shape)
            )

        # High-fidelity rendering additions: Auto Shape geometry type and borders
        try:
            if hasattr(shape, "auto_shape_type") and shape.auto_shape_type is not None:
                metadata["auto_shape_type"] = shape.auto_shape_type.name
        except Exception:
            pass


        try:
            if hasattr(shape, "line") and shape.line:
                border_color = self._get_safe_color_hex(shape.line.color)
                if border_color:
                    metadata["border_color"] = border_color
                if shape.line.width:
                    metadata["border_width"] = float(shape.line.width.pt)
        except Exception:
            pass


        return metadata


    def extract_table_data(self, shape):
        table = shape.table
        rows = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            rows.append(row_data)
        return rows

    def _extract_pptx_chart_data(self, shape) -> Optional[dict]:
        if not hasattr(shape, "chart") or shape.chart is None:
            return None
        try:
            chart = shape.chart
            title = ""
            if chart.has_title:
                title = chart.chart_title.text_frame.text

            series_data = []
            for s in chart.series:
                series_data.append({
                    "name": s.name,
                    "values": list(s.values) if hasattr(s, "values") else []
                })

            categories = []
            if len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, "categories"):
                    categories = [str(cat) for cat in plot.categories]

            chart_type_name = "unknown"
            if hasattr(chart, "chart_type"):
                chart_type_name = str(chart.chart_type)

            return {
                "title": title,
                "chart_type": chart_type_name,
                "series": series_data,
                "categories": categories
            }
        except Exception:
            return None

    def _extract_bullet_hierarchy(self, shape) -> list:
        items = []
        if not shape.has_text_frame:
            return items
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                items.append({"level": paragraph.level, "text": text})
        return items

    def _extract_connector_endpoints(self, shape) -> dict:
        endpoints: dict = {
            "begin_x": None,
            "begin_y": None,
            "end_x": None,
            "end_y": None,
        }
        # Prefer python-pptx Connector properties (reliable)
        if hasattr(shape, "begin_x"):
            try:
                endpoints["begin_x"] = float(shape.begin_x)
                endpoints["begin_y"] = float(shape.begin_y)
                endpoints["end_x"] = float(shape.end_x)
                endpoints["end_y"] = float(shape.end_y)
                return endpoints
            except Exception:
                pass
        # Fallback to XML parsing
        try:
            sp_element = shape._element
            xfrm = sp_element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
            )
            if xfrm is not None:
                off = xfrm.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}off"
                )
                ext = xfrm.find(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}ext"
                )
                if off is not None:
                    endpoints["begin_x"] = float(off.get("x", 0))
                    endpoints["begin_y"] = float(off.get("y", 0))
                if off is not None and ext is not None:
                    endpoints["end_x"] = (
                        float(off.get("x", 0)) + float(ext.get("cx", 0))
                    )
                    endpoints["end_y"] = (
                        float(off.get("y", 0)) + float(ext.get("cy", 0))
                    )
        except Exception:
            pass
        return endpoints

    def _safe_position(self, shape) -> tuple[float, float, float, float]:
        try:
            return (
                float(shape.left),
                float(shape.top),
                float(shape.width),
                float(shape.height),
            )
        except Exception:
            pass

        try:
            sp_element = shape._element
            xfrm = sp_element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
            )
            if xfrm is None:
                return 0.0, 0.0, 0.0, 0.0
            off = xfrm.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}off"
            )
            ext = xfrm.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}ext"
            )
            x = float(off.get("x", 0)) if off is not None else 0.0
            y = float(off.get("y", 0)) if off is not None else 0.0
            w = float(ext.get("cx", 0)) if ext is not None else 0.0
            h = float(ext.get("cy", 0)) if ext is not None else 0.0
            return x, y, w, h
        except Exception:
            return 0.0, 0.0, 0.0, 0.0


    def extract_table_as_markdown(self, shape) -> str:
        if not shape.has_table:
            return ""

        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
        body_lines = []
        for row in rows[1:]:
            body_lines.append("| " + " | ".join(row) + " |")

        parts = [header, separator] + body_lines
        return "\n".join(parts)

    def extract_table_as_list(self, shape):
        table_data = []

        for row in shape.table.rows:
            row_data = []

            for cell in row.cells:
                row_data.append(cell.text.strip())

            table_data.append(row_data)

        return table_data

    def extract_table_visual_metadata(self, shape):
        table = shape.table
        return {
        "row_count": len(table.rows),
        "column_count": len(table.columns),
        "column_widths": [col.width for col in table.columns],
        "row_heights": [row.height for row in table.rows],
    }

    def _get_safe_color_hex(self, color_obj) -> Optional[str]:
        """Safely extract hex color, avoiding ValueError on scheme colors."""
        if not color_obj:
            return None
        try:
            if color_obj.type == MSO_COLOR_TYPE.RGB and color_obj.rgb:
                return self.convert_rgb_to_hex(color_obj.rgb)
        except Exception:
            pass
        return None

    def _extract_text_style(self, shape) -> StyleModel:
        font_size: Optional[float] = None
        font_name: Optional[str] = None
        is_bold = False
        is_italic = False
        text_color = None
        background_color = None

        try:
            if shape.has_text_frame:
                paragraphs = shape.text_frame.paragraphs
                if paragraphs and paragraphs[0].runs:
                    first_run = paragraphs[0].runs[0]
                    font = first_run.font
                    if font.size:
                        font_size = float(font.size.pt)
                    if font.name:
                        font_name = font.name
                    if font.bold:
                        is_bold = True
                    if font.italic:
                        is_italic = True
                    text_color = self._get_safe_color_hex(font.color)
        except Exception:
            pass

        try:
            if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
                background_color = self.convert_rgb_to_hex(shape.fill.fore_color.rgb)
        except Exception:
            pass

        return StyleModel(
            font_size=font_size,
            font_name=font_name,
            bold=is_bold,
            italic=is_italic,
            text_color=text_color,
            background_color=background_color,
        )

    def sort_elements_by_reading_order(self, slide):
        """Sort elements top-to-bottom, then left-to-right."""
        return sorted(
            slide.elements,
            key=lambda e: (e.position.y, e.position.x),
        )

    @staticmethod
    def convert_rgb_to_hex(rgb_color: RGBColor) -> str:
        return "#{:02x}{:02x}{:02x}".format(
            rgb_color[0], rgb_color[1], rgb_color[2]
        )
